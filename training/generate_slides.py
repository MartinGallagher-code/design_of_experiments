#!/usr/bin/env python3
# Copyright (C) 2026 Martin J. Gallagher
# Licensed under the GNU General Public License v3.0 or later.
"""Generate PowerPoint slides for the DOE Helper Training Course.

Run: python generate_slides.py
Produces: slides/ directory with 8 .pptx files (one per module).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Brand colours
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
SECTION_BG = RGBColor(0xF5, 0xF7, 0xFA)
CODE_BG = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE0, 0x40, 0x40)

os.makedirs("slides", exist_ok=True)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_title_slide(prs, title, subtitle, module_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    if module_num is not None:
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(11), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = f"MODULE {module_num}"
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_BLUE
        p.font.bold = True
        p.font.name = "Calibri"

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Calibri"

    # Copyright footer
    tb3 = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11), Inches(0.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Copyright \u00a9 2026 Martin J. Gallagher. Licensed under GPL-3.0-or-later."
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p3.font.name = "Calibri"


def add_content_slide(prs, title, bullets, title_color=DARK_TEXT, bg_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)

    # Top accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"

    # Bullets
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11), Inches(5.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        # Support indent levels via leading "  "
        indent = 0
        text = bullet
        while text.startswith("  "):
            indent += 1
            text = text[2:]
        p.text = text
        p.font.size = Pt(22) if indent == 0 else Pt(19)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = indent
        if indent == 0:
            p.font.bold = False
    return slide


def add_code_slide(prs, title, code_text, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GREEN
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.font.name = "Calibri"

    # Code box
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.fill.background()

    tf2 = code_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)
    for i, line_text in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        p.font.name = "Consolas"
        p.space_after = Pt(2)

    if note:
        tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11), Inches(0.8))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.font.italic = True
        p3.font.name = "Calibri"
    return slide


def add_two_col_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.font.name = "Calibri"

    # Left column
    if left_title:
        tb_lt = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5))
        tf_lt = tb_lt.text_frame
        p_lt = tf_lt.paragraphs[0]
        p_lt.text = left_title
        p_lt.font.size = Pt(22)
        p_lt.font.color.rgb = ACCENT_BLUE
        p_lt.font.bold = True
        p_lt.font.name = "Calibri"

    y_start = Inches(2.0) if left_title else Inches(1.6)
    tb_l = slide.shapes.add_textbox(Inches(0.8), y_start, Inches(5.5), Inches(5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf_l.paragraphs[0] if i == 0 else tf_l.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(6)

    # Right column
    if right_title:
        tb_rt = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(5.5), Inches(0.5))
        tf_rt = tb_rt.text_frame
        p_rt = tf_rt.paragraphs[0]
        p_rt.text = right_title
        p_rt.font.size = Pt(22)
        p_rt.font.color.rgb = ACCENT_GREEN
        p_rt.font.bold = True
        p_rt.font.name = "Calibri"

    tb_r = slide.shapes.add_textbox(Inches(7), y_start, Inches(5.5), Inches(5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return slide


def add_exercise_slide(prs, title, instructions):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # Exercise badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(2.2), Inches(0.55)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    btf = badge.text_frame
    btf.paragraphs[0].text = "HANDS-ON EXERCISE"
    btf.paragraphs[0].font.size = Pt(14)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].font.name = "Calibri"
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tb = slide.shapes.add_textbox(Inches(3.3), Inches(0.35), Inches(9), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.font.name = "Calibri"

    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, inst in enumerate(instructions):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        indent = 0
        text = inst
        while text.startswith("  "):
            indent += 1
            text = text[2:]
        p.text = text
        p.font.size = Pt(20) if indent == 0 else Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = indent
    return slide


def add_key_point_slide(prs, title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"

    y = Inches(1.8)
    for point in points:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
        box.line.fill.background()
        btf = box.text_frame
        btf.margin_left = Inches(0.3)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = btf.paragraphs[0]
        p.text = point
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        y += Inches(1.05)
    return slide


def add_section_divider(prs, section_title, section_desc=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ACCENT_BLUE)

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    if section_desc:
        tb2 = slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(9), Inches(1))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = section_desc
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER


# ============================================================
# MODULE 1: Introduction to DOE
# ============================================================
def build_module_1():
    prs = new_prs()
    add_title_slide(prs,
        "Introduction to Design of Experiments",
        "Why systematic experimentation beats trial-and-error\n\nDOE Helper Training Course",
        module_num=1)

    add_content_slide(prs, "What You Will Learn", [
        "What Design of Experiments (DOE) is and why it matters",
        "The cost of unplanned experiments",
        "The One-Variable-At-a-Time (OVAT) trap",
        "Core DOE concepts: factors, levels, responses, effects",
        "How doe-helper automates the DOE workflow",
        "When to use DOE vs. other approaches",
    ])

    add_content_slide(prs, "What is Design of Experiments?", [
        "A systematic method for planning experiments",
        "  Change multiple factors simultaneously",
        "  Use statistical principles to extract maximum information",
        "  Minimize the number of experimental runs needed",
        "Pioneered by R.A. Fisher in the 1920s for agriculture",
        "Extended by Taguchi, Box, and others for industry",
        "Now applied to software, cloud, manufacturing, and beyond",
    ])

    add_two_col_slide(prs, "The Cost of Unplanned Experiments",
        [
            "Ad-hoc testing wastes resources",
            "Interactions between factors are missed",
            "Conclusions drawn from insufficient data",
            "Optimisation gets stuck in local optima",
            "Results are not reproducible",
        ],
        [
            "A web team ran 47 ad-hoc tests over 3 months",
            "Found: only 2 factors mattered",
            "A 2-factor full factorial needed just 4 runs",
            "DOE could have saved 43 runs and 11 weeks",
        ],
        left_title="The Problem",
        right_title="Real-World Example"
    )

    add_content_slide(prs, "The OVAT Trap", [
        "One-Variable-At-a-Time: change one factor, hold others constant",
        "Seems logical but has critical flaws:",
        "  Misses interactions between factors entirely",
        "  Requires more runs to learn the same amount",
        "  Can converge on wrong optimum",
        "Example: tuning a web server",
        "  OVAT: 3 factors x 3 levels = 9 runs, misses interactions",
        "  Full factorial: 2^3 = 8 runs, captures all interactions",
    ])

    add_content_slide(prs, "Core DOE Concepts", [
        "Factor: a variable you control (e.g. temperature, cache size)",
        "Level: specific value of a factor (e.g. low/high, 100/200/300)",
        "Response: what you measure (e.g. throughput, yield, cost)",
        "Main Effect: the impact of one factor on the response",
        "Interaction: when the effect of one factor depends on another",
        "Design Matrix: the table of all planned experimental runs",
        "Randomisation: running experiments in random order to avoid bias",
    ])

    add_content_slide(prs, "Effect Sparsity Principle", [
        "In most systems, only a few factors have large effects",
        "Most interactions are negligible",
        "This is the \"Pareto of experimentation\"",
        "Screening designs exploit this: test many factors cheaply",
        "Then focus detailed experiments on the vital few",
        "doe-helper supports this workflow: screen -> refine -> optimise",
    ])

    add_content_slide(prs, "The DOE Workflow", [
        "1. Define objectives: what do you want to learn or optimise?",
        "2. Select factors and levels: what can you control?",
        "3. Choose a design: how many runs can you afford?",
        "4. Generate the design matrix and run experiments",
        "5. Analyse results: which factors matter?",
        "6. Interpret and act: confirm, refine, or optimise",
        "doe-helper automates steps 3-6 and guides 1-2",
    ])

    add_code_slide(prs, "Meet doe-helper", """
$ pip install doehelper

$ doe init reactor --design full_factorial
  Created: reactor/config.json

$ doe generate reactor/
  Design matrix written to reactor/design.csv
  Runner script written to reactor/run.sh

$ doe analyze reactor/
  ANOVA table, Pareto chart, main-effects plots generated

$ doe optimize reactor/
  Recommended settings: temperature=180, pressure=3.5
""", "doe-helper provides a complete CLI workflow from design to optimisation")

    add_key_point_slide(prs, "Key Takeaways", [
        "DOE is a systematic, statistically rigorous way to experiment",
        "OVAT wastes runs and misses interactions",
        "Effect sparsity: most systems have only a few important factors",
        "doe-helper automates the entire DOE workflow via the CLI",
        "The workflow: define -> design -> run -> analyse -> optimise",
    ])

    add_exercise_slide(prs, "Exercise 1: Explore doe-helper", [
        "1. Install doe-helper: pip install doehelper",
        "2. Run: doe --help  and read the available commands",
        "3. Run: doe init coffee_brewing --design full_factorial",
        "4. Open coffee_brewing/config.json and examine the template",
        "5. Discussion: identify the factors, levels, and responses",
        "",
        "Bonus: Think of a problem in your own work where OVAT has been used.",
        "How many factors and levels would a proper DOE need?",
    ])

    prs.save("slides/Module_01_Introduction_to_DOE.pptx")
    print("  Module 1 saved")


# ============================================================
# MODULE 2: Getting Started with doe-helper
# ============================================================
def build_module_2():
    prs = new_prs()
    add_title_slide(prs,
        "Getting Started with doe-helper",
        "Installation, configuration, and your first experiment",
        module_num=2)

    add_content_slide(prs, "Learning Objectives", [
        "Install doe-helper and verify the setup",
        "Understand the config.json structure",
        "Define factors (continuous, categorical, ordinal)",
        "Define responses and their targets",
        "Configure design settings (randomisation, blocks, replicates)",
        "Use doe init with built-in templates",
    ])

    add_code_slide(prs, "Installation", """
# Install from PyPI
$ pip install doehelper

# Verify installation
$ doe --help

# Check version
$ doe --version
doehelper 0.1.0

# System requirements:
#   Python 3.10+
#   Works on Linux, macOS, Windows
#   Dependencies: numpy, pandas, scipy, matplotlib, pyDOE3, Jinja2
""")

    add_code_slide(prs, "Creating Your First Experiment", """
# Use a built-in template
$ doe init my_experiment --design full_factorial

# Or start from scratch: create config.json manually
$ mkdir my_experiment
$ cat > my_experiment/config.json << 'EOF'
{
  "factors": {
    "temperature": {"low": 150, "high": 200, "units": "C"},
    "pressure":    {"low": 1.0, "high": 5.0, "units": "bar"},
    "catalyst":    {"low": 0.5, "high": 1.5, "units": "g"}
  },
  "responses": {
    "yield": {"units": "%", "target": "maximize"}
  },
  "design": {"type": "full_factorial"}
}
EOF
""")

    add_content_slide(prs, "Factor Types", [
        "Continuous: numeric range with low/high values",
        "  Example: temperature (150-200), pressure (1.0-5.0)",
        "Categorical: discrete named levels",
        "  Example: material (\"steel\", \"aluminum\", \"titanium\")",
        "Ordinal: ordered categories with inherent ranking",
        "  Example: quality (\"low\", \"medium\", \"high\")",
        "",
        "Coded variables: factors are internally scaled to -1 / +1",
        "  This ensures all effects are directly comparable",
    ])

    add_code_slide(prs, "Factor Configuration Examples", """
"factors": {
  "temperature": {
    "low": 150, "high": 200, "units": "C"
  },
  "material": {
    "type": "categorical",
    "levels": ["steel", "aluminum", "titanium"]
  },
  "quality_grade": {
    "type": "ordinal",
    "levels": ["low", "medium", "high"]
  }
}
""", "Continuous factors use low/high; categorical and ordinal use levels list")

    add_content_slide(prs, "Response Configuration", [
        "Each response has a name, units, and a target direction",
        "Target options:",
        "  \"maximize\" - higher is better (yield, throughput)",
        "  \"minimize\" - lower is better (cost, latency, defects)",
        "  A specific numeric value - hit a target (pH = 7.0)",
        "Multiple responses are supported for multi-objective optimisation",
        "  doe-helper uses desirability functions to balance trade-offs",
    ])

    add_code_slide(prs, "Response Configuration Examples", """
"responses": {
  "yield": {
    "units": "%",
    "target": "maximize"
  },
  "cost": {
    "units": "USD",
    "target": "minimize"
  },
  "purity": {
    "units": "%",
    "target": 99.5
  }
}
""")

    add_content_slide(prs, "Design Settings", [
        "The \"design\" section controls how the experiment is structured:",
        "  type: design type (full_factorial, ccd, box_behnken, etc.)",
        "  randomize: true/false - randomise run order (default: true)",
        "  replicates: number of times to repeat each run (default: 1)",
        "  center_points: add center points for curvature detection",
        "  blocks: number of blocks for blocking out nuisance variables",
        "  resolution: for fractional factorials (III, IV, V)",
    ])

    add_code_slide(prs, "Full config.json Example", """
{
  "factors": {
    "temperature": {"low": 150, "high": 200, "units": "C"},
    "pressure":    {"low": 1.0, "high": 5.0, "units": "bar"},
    "catalyst":    {"low": 0.5, "high": 1.5, "units": "g"}
  },
  "responses": {
    "yield":  {"units": "%", "target": "maximize"},
    "purity": {"units": "%", "target": "maximize"}
  },
  "design": {
    "type": "full_factorial",
    "randomize": true,
    "center_points": 3
  }
}
""", "This config produces a 2^3 = 8 runs + 3 center points = 11 total runs")

    add_content_slide(prs, "Built-in Templates (doe init)", [
        "doe-helper includes 221 worked use-case templates",
        "Categories include:",
        "  Cloud/DevOps: Kubernetes, CI/CD, database tuning",
        "  Manufacturing: 3D printing, injection molding, chemical processes",
        "  Food science: coffee brewing, bread baking, fermentation",
        "  IoT/Electronics: sensor calibration, motor control",
        "  Sports/Health: training programs, sleep optimisation",
        "Templates provide config.json + simulation scripts for practice",
    ])

    add_code_slide(prs, "Using doe info", """
$ doe info my_experiment/

Design Summary
  Type:           full_factorial
  Factors:        3
  Runs:           11 (8 factorial + 3 center)
  Replicates:     1
  Randomized:     yes

Design Evaluation
  D-efficiency:   100.0%
  A-efficiency:   100.0%
  G-efficiency:   100.0%
""", "Use doe info to verify your design before running experiments")

    add_exercise_slide(prs, "Exercise 2: Build a Configuration", [
        "Scenario: You're optimising a web application's performance.",
        "",
        "1. Create a directory: mkdir webapp_perf",
        "2. Create config.json with these factors:",
        "  cache_size: 64 to 512 (MB)",
        "  thread_count: 4 to 32",
        "  compression: categorical [\"none\", \"gzip\", \"brotli\"]",
        "3. Add two responses:",
        "  response_time_ms (minimize) and throughput_rps (maximize)",
        "4. Set the design type to full_factorial",
        "5. Run: doe info webapp_perf/  to verify your configuration",
        "6. Run: doe generate webapp_perf/  to see the design matrix",
    ])

    prs.save("slides/Module_02_Getting_Started.pptx")
    print("  Module 2 saved")


# ============================================================
# MODULE 3: Full Factorial Designs
# ============================================================
def build_module_3():
    prs = new_prs()
    add_title_slide(prs,
        "Full Factorial Designs",
        "Understanding and running complete two-level experiments",
        module_num=3)

    add_content_slide(prs, "Learning Objectives", [
        "Understand full factorial (2^k) designs",
        "Calculate the number of runs needed",
        "Interpret main effects and interactions",
        "Generate and run a full factorial with doe-helper",
        "Read ANOVA tables and Pareto charts",
        "Know when full factorial is the right choice",
    ])

    add_content_slide(prs, "What is a Full Factorial Design?", [
        "Tests every combination of factor levels",
        "For k factors at 2 levels: 2^k runs",
        "  2 factors: 4 runs",
        "  3 factors: 8 runs",
        "  4 factors: 16 runs",
        "  5 factors: 32 runs",
        "Estimates ALL main effects and ALL interactions",
        "The \"gold standard\" of experimental designs",
        "Practical limit: ~5-6 factors (32-64 runs)",
    ])

    add_two_col_slide(prs, "2^2 Full Factorial Example",
        [
            "Factor A: Temperature (150, 200)",
            "Factor B: Pressure (1, 5)",
            "",
            "Run 1: A=low,  B=low   -> 62%",
            "Run 2: A=high, B=low   -> 74%",
            "Run 3: A=low,  B=high  -> 68%",
            "Run 4: A=high, B=high  -> 91%",
        ],
        [
            "Main effect A = [(74+91)-(62+68)]/2",
            "  = 17.5% (Temperature matters)",
            "",
            "Main effect B = [(68+91)-(62+74)]/2",
            "  = 11.5% (Pressure matters)",
            "",
            "Interaction AB = [(62+91)-(74+68)]/2",
            "  = 5.5% (Synergy!)",
        ],
        left_title="Design Matrix & Results",
        right_title="Effect Calculations"
    )

    add_content_slide(prs, "Why Interactions Matter", [
        "An interaction means the effect of one factor depends on another",
        "OVAT cannot detect interactions at all",
        "Example: Temperature alone increases yield by 17.5%",
        "  But at high pressure, temperature boost is even larger",
        "  The combination is more powerful than either alone",
        "Interactions can be antagonistic too (cancel each other out)",
        "Full factorials capture all of these relationships",
    ])

    add_code_slide(prs, "Running a Full Factorial with doe-helper", """
# 1. Configure
$ cat > seal_strength/config.json << 'EOF'
{
  "factors": {
    "temperature": {"low": 225, "high": 285, "units": "F"},
    "pressure":    {"low": 40,  "high": 90,  "units": "psi"},
    "dwell_time":  {"low": 1.0, "high": 3.5, "units": "sec"}
  },
  "responses": {
    "seal_strength": {"units": "g/in", "target": "maximize"}
  },
  "design": {"type": "full_factorial", "center_points": 3}
}
EOF

# 2. Generate the design matrix
$ doe generate seal_strength/

# 3. Run experiments and record results
$ doe record seal_strength/

# 4. Analyse
$ doe analyze seal_strength/
""")

    add_content_slide(prs, "Reading the ANOVA Table", [
        "ANOVA = Analysis of Variance",
        "Key columns:",
        "  Source: factor or interaction name",
        "  SS: Sum of Squares (variability explained)",
        "  DF: Degrees of Freedom",
        "  MS: Mean Square (SS / DF)",
        "  F-value: test statistic (higher = stronger effect)",
        "  p-value: probability effect is due to chance",
        "Rule of thumb: p < 0.05 means the factor is significant",
    ])

    add_content_slide(prs, "Reading the Pareto Chart", [
        "Bar chart of absolute effect magnitudes, sorted largest to smallest",
        "Dashed reference line shows significance threshold",
        "Bars above the line are statistically significant effects",
        "Quickly identifies the \"vital few\" factors",
        "doe-helper generates this automatically with doe analyze",
    ])

    add_content_slide(prs, "Center Points", [
        "Extra runs at the midpoint of all factors",
        "Purpose: detect curvature (non-linear effects)",
        "If center-point average differs from predicted linear value:",
        "  The true relationship is curved",
        "  You may need a response surface design",
        "Cheap insurance: 3-5 center points cost little",
        "doe-helper: set center_points in config.json",
    ])

    add_key_point_slide(prs, "Key Takeaways", [
        "Full factorial = every combination, captures all effects",
        "2^k runs: practical for up to ~5-6 factors",
        "Interactions are why DOE beats OVAT",
        "Center points detect curvature cheaply",
        "doe-helper automates generation, recording, and analysis",
    ])

    add_exercise_slide(prs, "Exercise 3: Seal Strength Experiment", [
        "1. Create the seal_strength experiment from the slide",
        "  mkdir seal_strength && create config.json as shown",
        "2. Run: doe generate seal_strength/",
        "3. Run: doe info seal_strength/  -- how many runs?",
        "4. Use the simulation script:  bash seal_strength/run.sh",
        "5. Run: doe analyze seal_strength/",
        "6. Answer these questions:",
        "  Which factor has the largest main effect?",
        "  Are there any significant interactions?",
        "  Is there evidence of curvature from center points?",
        "7. Run: doe report seal_strength/  -- open the HTML report",
    ])

    prs.save("slides/Module_03_Full_Factorial.pptx")
    print("  Module 3 saved")


# ============================================================
# MODULE 4: Fractional Factorial & Screening
# ============================================================
def build_module_4():
    prs = new_prs()
    add_title_slide(prs,
        "Fractional Factorial & Screening Designs",
        "Testing many factors efficiently when resources are limited",
        module_num=4)

    add_content_slide(prs, "Learning Objectives", [
        "Understand why fractional factorials exist",
        "Learn about resolution and aliasing",
        "Use Plackett-Burman and Definitive Screening designs",
        "Choose the right screening design with doe-helper",
        "Interpret confounded effects",
        "Plan a follow-up experiment after screening",
    ])

    add_content_slide(prs, "The Problem: Too Many Factors", [
        "Full factorial with 7 factors = 128 runs",
        "With 10 factors = 1,024 runs",
        "With 15 factors = 32,768 runs!",
        "In practice, many experiments start with 6-20 factors",
        "Effect sparsity says most factors won't matter",
        "Solution: run a fraction of the full factorial",
    ])

    add_content_slide(prs, "Fractional Factorial Designs", [
        "Run 2^(k-p) instead of 2^k runs",
        "  2^(7-4) = 8 runs instead of 128 for 7 factors",
        "  2^(6-2) = 16 runs instead of 64 for 6 factors",
        "Trade-off: some effects become aliased (confounded)",
        "Resolution tells you what is confounded with what:",
        "  Resolution III: main effects aliased with 2-factor interactions",
        "  Resolution IV: main effects clear, 2FI aliased with each other",
        "  Resolution V: main effects and 2FI all clear",
    ])

    add_code_slide(prs, "Fractional Factorial with doe-helper", """
{
  "factors": {
    "A": {"low": -1, "high": 1},
    "B": {"low": -1, "high": 1},
    "C": {"low": -1, "high": 1},
    "D": {"low": -1, "high": 1},
    "E": {"low": -1, "high": 1},
    "F": {"low": -1, "high": 1},
    "G": {"low": -1, "high": 1}
  },
  "responses": {
    "performance": {"units": "score", "target": "maximize"}
  },
  "design": {"type": "fractional_factorial"}
}

# doe-helper automatically selects the appropriate resolution
$ doe generate screening_exp/
$ doe info screening_exp/  # shows aliasing structure
""")

    add_content_slide(prs, "Plackett-Burman Designs", [
        "Resolution III screening designs",
        "Run counts in multiples of 4: 12, 20, 24, 28, 36...",
        "Screen up to N-1 factors in N runs",
        "  12 runs -> up to 11 factors",
        "  20 runs -> up to 19 factors",
        "Very efficient for identifying the vital few factors",
        "Main effects are confounded with 2-factor interactions",
        "Use: initial screening when you have many candidate factors",
    ])

    add_code_slide(prs, "Plackett-Burman with doe-helper", """
{
  "factors": {
    "cache_ttl":      {"low": 60,   "high": 3600,  "units": "sec"},
    "max_connections": {"low": 10,   "high": 100},
    "thread_pool":    {"low": 4,    "high": 64},
    "buffer_size":    {"low": 1024, "high": 65536, "units": "bytes"},
    "retry_count":    {"low": 0,    "high": 5},
    "timeout":        {"low": 1,    "high": 30,    "units": "sec"},
    "batch_size":     {"low": 10,   "high": 1000}
  },
  "responses": {
    "latency_p99": {"units": "ms", "target": "minimize"}
  },
  "design": {"type": "plackett_burman"}
}

# Only 12 runs to screen 7 factors!
$ doe generate api_screening/
""")

    add_content_slide(prs, "Definitive Screening Designs", [
        "Modern alternative to Plackett-Burman (Jones & Nachtsheim, 2011)",
        "Three-level designs: estimate some quadratic effects",
        "2k+1 runs for k continuous factors",
        "Advantages over Plackett-Burman:",
        "  Main effects not aliased with two-factor interactions",
        "  Can detect curvature without extra runs",
        "  Supported for 3+ continuous factors",
        "doe-helper: set \"type\": \"definitive_screening\"",
    ])

    add_content_slide(prs, "Choosing a Screening Design", [
        "Plackett-Burman: maximum economy, many factors, need follow-up",
        "Definitive Screening: slightly more runs, cleaner estimates",
        "Fractional Factorial: when you need specific resolution",
        "",
        "Decision guide:",
        "  > 6 factors, tight budget -> Plackett-Burman or DSD",
        "  4-6 factors, moderate budget -> Fractional Factorial (Res IV+)",
        "  <= 5 factors, can afford it -> Full Factorial",
        "doe-helper handles the design selection automatically",
    ])

    add_two_col_slide(prs, "The Screening Workflow",
        [
            "Phase 1: Screen",
            "  Run PB or DSD with all candidate factors",
            "  Identify 3-5 significant factors",
            "  Drop non-significant factors",
            "",
            "Phase 2: Characterise",
            "  Full factorial on surviving factors",
            "  Estimate all interactions",
            "  Check for curvature",
        ],
        [
            "Phase 3: Optimise",
            "  Response surface design (CCD, BBD)",
            "  Model the response surface",
            "  Find the optimum settings",
            "",
            "doe-helper supports all phases:",
            "  doe generate -> doe analyze",
            "  doe augment for fold-overs",
            "  doe optimize for final settings",
        ],
        left_title="Screen & Characterise",
        right_title="Optimise"
    )

    add_code_slide(prs, "Fold-Over with doe augment", """
# After screening, if you need to de-alias effects:
$ doe augment api_screening/ --method fold-over

# This mirrors the original design to break aliases
# Doubles the run count but separates confounded effects

# Other augmentation methods:
$ doe augment experiment/ --method center-points --count 5
$ doe augment experiment/ --method star-points  # for CCD
""", "doe augment extends an existing design without starting over")

    add_key_point_slide(prs, "Key Takeaways", [
        "Fractional factorials trade runs for information (aliasing)",
        "Resolution tells you what's confounded with what",
        "Plackett-Burman: maximum screening efficiency",
        "Definitive Screening: modern, cleaner main-effect estimates",
        "doe-helper automates design selection and fold-overs",
    ])

    add_exercise_slide(prs, "Exercise 4: Screening a Microservice", [
        "Scenario: 8 factors that might affect API latency.",
        "",
        "1. Create config.json with 8 continuous factors of your choice",
        "  (e.g. connection_pool, timeout, retries, batch_size, ...)",
        "2. Set design type to \"plackett_burman\"",
        "3. Run: doe generate microservice/",
        "  How many runs does the design require?",
        "4. Run: doe info microservice/",
        "5. Compare: change design type to \"definitive_screening\"",
        "  How many runs now? What's different?",
        "6. Run the simulation and analyse results",
        "7. Which factors appear significant? Plan a follow-up design.",
    ])

    prs.save("slides/Module_04_Screening_Designs.pptx")
    print("  Module 4 saved")


# ============================================================
# MODULE 5: Response Surface Designs
# ============================================================
def build_module_5():
    prs = new_prs()
    add_title_slide(prs,
        "Response Surface Designs",
        "Modeling curved relationships and finding optimal settings",
        module_num=5)

    add_content_slide(prs, "Learning Objectives", [
        "Understand when you need response surface methodology (RSM)",
        "Learn Central Composite Design (CCD) and Box-Behnken (BBD)",
        "Use Latin Hypercube Sampling for space-filling",
        "Generate RSM designs with doe-helper",
        "Interpret 3D surface plots and contour plots",
        "Use doe optimize to find optimal factor settings",
    ])

    add_content_slide(prs, "When to Move Beyond Screening", [
        "Screening identified the vital few factors (typically 2-4)",
        "Center points showed evidence of curvature",
        "You need to find the optimum, not just know what matters",
        "Response Surface Methodology (RSM) fits a quadratic model:",
        "  y = b0 + b1*x1 + b2*x2 + b12*x1*x2 + b11*x1^2 + b22*x2^2",
        "This captures curves, ridges, saddle points, and optima",
    ])

    add_content_slide(prs, "Central Composite Design (CCD)", [
        "The most popular RSM design",
        "Three components:",
        "  Factorial points: 2^k corner points",
        "  Star (axial) points: 2k points along each axis",
        "  Center points: 3-6 replicates at the center",
        "For 3 factors: 8 + 6 + 6 = 20 runs",
        "Can be run sequentially: factorial first, then augment with star points",
        "Star point distance (alpha) controls rotatability",
    ])

    add_content_slide(prs, "Box-Behnken Design", [
        "Alternative to CCD -- fewer runs for 3+ factors",
        "Does not include corner points (extreme combinations)",
        "Good when extremes are costly or dangerous",
        "For 3 factors: 15 runs (vs. 20 for CCD)",
        "For 4 factors: 27 runs (vs. 30 for CCD)",
        "Slightly less information than CCD but often sufficient",
        "doe-helper: set \"type\": \"box_behnken\"",
    ])

    add_code_slide(prs, "CCD with doe-helper", """
{
  "factors": {
    "temperature": {"low": 150, "high": 200, "units": "C"},
    "pressure":    {"low": 1.0, "high": 5.0, "units": "bar"},
    "catalyst":    {"low": 0.5, "high": 1.5, "units": "g"}
  },
  "responses": {
    "yield":  {"units": "%", "target": "maximize"},
    "purity": {"units": "%", "target": "maximize"}
  },
  "design": {
    "type": "ccd",
    "center_points": 6
  }
}

$ doe generate reactor_rsm/
$ doe info reactor_rsm/
  Type: ccd | Runs: 20 | Factors: 3
""")

    add_content_slide(prs, "Latin Hypercube Sampling (LHS)", [
        "Space-filling design: samples spread evenly across the space",
        "Not model-based -- does not assume a specific model form",
        "Good for computer experiments and simulations",
        "User specifies the number of runs",
        "Each factor divided into N equal intervals, one sample per interval",
        "Complements DOE when the response surface is unknown",
        "doe-helper: set \"type\": \"latin_hypercube\", \"runs\": N",
    ])

    add_content_slide(prs, "Interpreting Surface Plots", [
        "3D surface plots show response as a function of two factors",
        "Contour plots show the same information as a 2D map",
        "Look for:",
        "  Peaks/valleys: where the optimum lies",
        "  Ridges: elongated regions of near-optimal response",
        "  Saddle points: optimum in one direction, pessimum in another",
        "doe-helper generates these automatically in reports",
        "  doe report experiment/ produces an interactive HTML report",
    ])

    add_code_slide(prs, "Optimization with doe-helper", """
# After recording results and running analysis:
$ doe optimize reactor_rsm/

Optimization Results
  Method: L-BFGS-B with 20 random restarts

  Optimal Settings:
    temperature: 182.3 C
    pressure:    3.8 bar
    catalyst:    1.2 g

  Predicted Responses:
    yield:  94.2%  (target: maximize)
    purity: 98.7%  (target: maximize)

  Desirability: 0.93
""", "doe optimize uses scipy.optimize with multi-start for global search")

    add_content_slide(prs, "Steepest Ascent / Descent", [
        "Sequential experimentation strategy",
        "Start with a factorial/screening design near current settings",
        "Fit a linear model to identify the steepest ascent direction",
        "Run experiments along this path until response stops improving",
        "Then run a new RSM design at the plateau",
        "doe-helper computes the steepest ascent path automatically",
        "Efficient when you're far from the optimum",
    ])

    add_content_slide(prs, "Design Type Selection Guide", [
        "2-3 factors, need optimum -> CCD or Box-Behnken",
        "3-4 factors, avoid extremes -> Box-Behnken",
        "Unknown model, computer sim -> Latin Hypercube",
        "Augmenting a factorial -> add star points (doe augment)",
        "Mixture problems -> Simplex-Lattice or Simplex-Centroid",
        "Constrained factor space -> D-Optimal",
        "doe-helper supports all 11 design types",
    ])

    add_key_point_slide(prs, "Key Takeaways", [
        "RSM fits quadratic models to find optima",
        "CCD is the workhorse; Box-Behnken avoids extremes",
        "LHS for space-filling when the model is unknown",
        "doe optimize uses L-BFGS-B with multi-start",
        "Use steepest ascent when far from the optimum",
    ])

    add_exercise_slide(prs, "Exercise 5: Response Surface Optimisation", [
        "1. Start from the reactor_optimization use case:",
        "  doe init reactor_optimization --design ccd",
        "2. Run: doe generate reactor_optimization/",
        "3. Examine the design matrix: how many runs? What's the structure?",
        "4. Run the simulation: bash reactor_optimization/run.sh",
        "5. Run: doe analyze reactor_optimization/",
        "  Look at the response surface plots",
        "6. Run: doe optimize reactor_optimization/",
        "  What are the recommended optimal settings?",
        "7. Now try Box-Behnken: change type to \"box_behnken\" and repeat",
        "  Compare the run count and optimal settings",
    ])

    prs.save("slides/Module_05_Response_Surface.pptx")
    print("  Module 5 saved")


# ============================================================
# MODULE 6: Analysis & Interpretation
# ============================================================
def build_module_6():
    prs = new_prs()
    add_title_slide(prs,
        "Analysis and Interpretation",
        "Turning experimental data into actionable insights",
        module_num=6)

    add_content_slide(prs, "Learning Objectives", [
        "Read and interpret ANOVA tables confidently",
        "Use Pareto charts and main-effects plots",
        "Understand residual analysis and model diagnostics",
        "Detect and handle model inadequacy",
        "Use doe analyze for automated analysis",
        "Generate comprehensive HTML reports with doe report",
    ])

    add_content_slide(prs, "The doe analyze Output", [
        "doe analyze produces a complete analysis package:",
        "  ANOVA table with F-tests and p-values",
        "  Effect estimates with confidence intervals",
        "  Pareto chart of effect magnitudes",
        "  Main effects plots",
        "  Interaction plots",
        "  Normal probability plot of effects",
        "  Residual diagnostics (4-panel plot)",
        "  Model summary (R-squared, adjusted R-squared)",
    ])

    add_content_slide(prs, "ANOVA Table Deep Dive", [
        "Source: name of the factor, interaction, or error term",
        "Sum of Squares (SS): total variability explained by this source",
        "  % Contribution = SS_source / SS_total x 100",
        "Degrees of Freedom (DF): parameters used",
        "Mean Square (MS): SS / DF",
        "F-value: MS_source / MS_error (signal-to-noise ratio)",
        "p-value: probability of seeing this F by chance",
        "  p < 0.05: significant | p < 0.01: highly significant",
    ])

    add_content_slide(prs, "Pareto Chart of Effects", [
        "Horizontal bar chart, absolute effects sorted large to small",
        "Dashed line = significance threshold",
        "  For replicated designs: based on pooled error estimate",
        "  For unreplicated: uses Lenth's pseudo-standard-error",
        "Focus on bars that cross the threshold",
        "The \"vital few\" factors are immediately obvious",
        "doe-helper labels each bar with the factor/interaction name",
    ])

    add_content_slide(prs, "Main Effects and Interaction Plots", [
        "Main Effects Plot:",
        "  Shows average response at each level of each factor",
        "  Steep slope = large effect; flat = no effect",
        "  Quick visual scan of which factors matter",
        "",
        "Interaction Plot:",
        "  Shows response for each combination of two factors",
        "  Parallel lines = no interaction",
        "  Non-parallel (crossing) lines = interaction present",
    ])

    add_content_slide(prs, "Normal Probability Plot of Effects", [
        "Plots effect estimates against expected normal quantiles",
        "Inactive effects fall on a straight line through zero",
        "Active effects deviate from the line",
        "Particularly useful for unreplicated designs",
        "  No error estimate available -> use this visual method",
        "doe-helper generates both normal and half-normal plots",
    ])

    add_content_slide(prs, "Residual Diagnostics", [
        "Residuals = observed - predicted values",
        "Four-panel diagnostic plot from doe analyze:",
        "  1. Residuals vs. Fitted: check for patterns (should be random)",
        "  2. Normal Q-Q plot: residuals should follow a normal distribution",
        "  3. Residuals vs. Run order: check for time trends",
        "  4. Residuals histogram: should be roughly bell-shaped",
        "Patterns in residuals signal model problems:",
        "  Funnel shape -> non-constant variance (transform response)",
        "  Curve -> missing quadratic terms (need RSM)",
    ])

    add_content_slide(prs, "Model Adequacy Statistics", [
        "R-squared: fraction of variability explained (0 to 1)",
        "  R^2 > 0.9 is generally good",
        "Adjusted R-squared: penalised for number of terms",
        "  Should be close to R^2 (large gap = overfitting)",
        "Predicted R-squared (PRESS-based): cross-validation metric",
        "  Should agree with Adj R^2 within ~0.2",
        "Lack-of-Fit test: is the model adequate?",
        "  Significant LOF -> model is missing important terms",
    ])

    add_code_slide(prs, "doe analyze in Practice", """
$ doe analyze reactor_rsm/

ANOVA Table
  Source        SS       DF    MS       F       p-value   Sig
  temperature   1240.5   1    1240.5   45.2    0.0001    ***
  pressure       890.3   1     890.3   32.4    0.0005    ***
  catalyst       45.2    1      45.2    1.6    0.2345
  temp*press     320.1   1     320.1   11.7    0.0089    **
  temp^2         156.7   1     156.7    5.7    0.0412    *
  ...

Model: R^2 = 0.967   Adj R^2 = 0.945   Pred R^2 = 0.912
""")

    add_code_slide(prs, "Generating HTML Reports", """
$ doe report reactor_rsm/

Report generated: reactor_rsm/report.html

# The report includes:
#   - Executive summary with key findings
#   - Interactive plots (zoom, hover, pan)
#   - Full ANOVA table
#   - Effect estimates and confidence intervals
#   - Residual diagnostics
#   - Response surface plots (for RSM designs)
#   - Optimal settings (if doe optimize was run)
#   - Self-contained single HTML file -- easy to share
""")

    add_key_point_slide(prs, "Key Takeaways", [
        "doe analyze automates the entire analysis pipeline",
        "Pareto chart: quick visual ID of significant factors",
        "Always check residual plots for model adequacy",
        "R^2, Adj-R^2, and Pred-R^2 should be close and high",
        "doe report creates shareable HTML reports with all results",
    ])

    add_exercise_slide(prs, "Exercise 6: Analyse and Interpret", [
        "Use the reactor_optimization experiment from Exercise 5.",
        "",
        "1. Run: doe analyze reactor_optimization/",
        "2. Examine the ANOVA table:",
        "  Which factors are significant at p < 0.05?",
        "  What is the % contribution of each factor?",
        "3. Look at the Pareto chart: do the results agree?",
        "4. Check the residual diagnostics:",
        "  Are there patterns in residuals vs. fitted values?",
        "  Do residuals follow a normal distribution?",
        "5. Run: doe report reactor_optimization/",
        "  Open the HTML report and explore the interactive plots",
        "6. Write a brief summary: What factors matter? Any interactions?",
    ])

    prs.save("slides/Module_06_Analysis.pptx")
    print("  Module 6 saved")


# ============================================================
# MODULE 7: Multi-Response Optimisation & Advanced Topics
# ============================================================
def build_module_7():
    prs = new_prs()
    add_title_slide(prs,
        "Multi-Response Optimisation\nand Advanced Topics",
        "Balancing competing objectives and extending your designs",
        module_num=7)

    add_content_slide(prs, "Learning Objectives", [
        "Optimise for multiple responses simultaneously",
        "Understand desirability functions",
        "Use design augmentation strategies",
        "Apply blocking and randomisation correctly",
        "Use Taguchi designs for robust parameter design",
        "Handle mixture experiments",
    ])

    add_content_slide(prs, "Multi-Response Optimisation", [
        "Most real experiments have multiple responses",
        "  Maximise yield AND minimise cost",
        "  Minimise latency AND maximise throughput",
        "Responses often conflict: improving one degrades another",
        "Need a systematic way to balance trade-offs",
        "Derringer-Suich desirability function:",
        "  Transform each response to a 0-1 desirability scale",
        "  Overall desirability D = geometric mean of individual d_i",
    ])

    add_code_slide(prs, "Multi-Response Config", """
{
  "factors": {
    "temperature": {"low": 150, "high": 200, "units": "C"},
    "pressure":    {"low": 1.0, "high": 5.0, "units": "bar"},
    "catalyst":    {"low": 0.5, "high": 1.5, "units": "g"}
  },
  "responses": {
    "yield":  {"units": "%",   "target": "maximize"},
    "cost":   {"units": "USD", "target": "minimize"},
    "purity": {"units": "%",   "target": 99.5}
  },
  "design": {"type": "ccd"}
}

$ doe optimize reactor_multi/
  Overall desirability: 0.87
  yield=91.2%, cost=$42.30, purity=99.3%
""", "doe optimize balances all responses via desirability functions")

    add_content_slide(prs, "Design Augmentation", [
        "Often you want to extend an existing design",
        "doe augment supports several strategies:",
        "  Fold-over: mirror the design to break aliases",
        "  Star points: add axial points for RSM",
        "  Center points: add runs at the center",
        "Benefits:",
        "  Leverages data you already have",
        "  Sequentially builds up information",
        "  No need to start over from scratch",
    ])

    add_content_slide(prs, "Blocking", [
        "Blocks account for known nuisance variables",
        "  Different batches of raw material",
        "  Different days or operators",
        "  Different machines",
        "Block effects are estimated and removed from the analysis",
        "doe-helper: set \"blocks\" in the design section",
        "The design is constructed so blocks are orthogonal to treatments",
    ])

    add_content_slide(prs, "Taguchi Designs", [
        "Focus on robustness: find settings insensitive to noise",
        "Inner array: controllable factors",
        "Outer array: noise factors (things you can't control)",
        "Signal-to-noise (S/N) ratio as the response",
        "doe-helper: set \"type\": \"taguchi\"",
        "Useful for manufacturing where environmental variation matters",
    ])

    add_content_slide(prs, "Mixture Experiments", [
        "When factors are proportions that must sum to 1 (or 100%)",
        "  Concrete: cement, water, aggregate, admixture",
        "  Alloy: proportions of different metals",
        "  Formulation: ingredients in a product",
        "Standard designs don't work (can't vary independently)",
        "doe-helper supports:",
        "  Simplex-Lattice: evenly-spaced mixture points",
        "  Simplex-Centroid: vertices, edge midpoints, center",
    ])

    add_content_slide(prs, "D-Optimal Designs", [
        "Computer-generated designs for non-standard situations:",
        "  Irregular factor space (constraints)",
        "  Mix of continuous and categorical factors",
        "  Limited run budget",
        "Maximises the determinant of X'X (D-optimality)",
        "Uses the Fedorov exchange algorithm",
        "doe-helper: set \"type\": \"d_optimal\", \"runs\": N",
        "Flexible but requires specifying the number of runs",
    ])

    add_code_slide(prs, "Power Analysis with doe-helper", """
$ doe power experiment/

Power Analysis
  Design:          full_factorial (2^3)
  Runs:            8
  Error estimate:  2.5 (from prior data or estimate)

  Detectable Effect Sizes (at 80% power, alpha=0.05):
    Main effects:   3.2 units
    2-factor int.:  3.2 units

# Or specify a target effect size:
$ doe power experiment/ --effect-size 5.0
  Power: 0.95 for main effects
  Power: 0.95 for 2-factor interactions
""", "Use power analysis to ensure your design can detect meaningful effects")

    add_key_point_slide(prs, "Key Takeaways", [
        "Multi-response: desirability functions balance trade-offs",
        "doe augment extends designs without starting over",
        "Blocking removes nuisance variation systematically",
        "Taguchi designs optimise for robustness to noise",
        "D-Optimal handles constraints and irregular spaces",
    ])

    add_exercise_slide(prs, "Exercise 7: Multi-Response Optimisation", [
        "1. Modify the reactor experiment to have 3 responses:",
        "  yield (maximize), cost (minimize), purity (target: 99.5)",
        "2. Run: doe generate reactor_multi/",
        "3. Run the simulation and collect results",
        "4. Run: doe analyze reactor_multi/",
        "5. Run: doe optimize reactor_multi/",
        "  What is the overall desirability score?",
        "  How do the optimal settings differ from single-response?",
        "6. Run: doe power reactor_multi/",
        "  Is the design adequately powered for all responses?",
        "7. Try: doe augment reactor_multi/ --method center-points --count 4",
        "  How does this change the power analysis?",
    ])

    prs.save("slides/Module_07_Advanced_Topics.pptx")
    print("  Module 7 saved")


# ============================================================
# MODULE 8: Capstone Project
# ============================================================
def build_module_8():
    prs = new_prs()
    add_title_slide(prs,
        "Capstone Project",
        "Apply the full DOE workflow to a real-world problem",
        module_num=8)

    add_content_slide(prs, "Capstone Overview", [
        "Apply everything you've learned in one complete project",
        "Choose from three scenario options or propose your own",
        "Complete the full DOE workflow:",
        "  1. Problem definition and factor selection",
        "  2. Screening design",
        "  3. Follow-up design (factorial or RSM)",
        "  4. Analysis and interpretation",
        "  5. Optimisation",
        "  6. Final report",
    ])

    add_content_slide(prs, "Option A: Cloud Infrastructure Optimisation", [
        "Scenario: optimise a Kubernetes-deployed microservice",
        "10 candidate factors: CPU limit, memory limit, replicas,",
        "  connection pool, thread count, cache TTL, batch size,",
        "  retry limit, timeout, log level",
        "3 responses: p99 latency (min), throughput (max), cost (min)",
        "Budget: 40 total experimental runs",
        "",
        "Use doe-helper to design, simulate, analyse, and optimise.",
    ])

    add_content_slide(prs, "Option B: Manufacturing Process", [
        "Scenario: optimise 3D printing parameters",
        "8 candidate factors: layer height, print speed, nozzle temp,",
        "  bed temp, infill %, infill pattern, retraction distance,",
        "  cooling fan speed",
        "3 responses: tensile strength (max), print time (min), surface quality (max)",
        "Budget: 30 total experimental runs",
        "",
        "Use doe-helper to design, simulate, analyse, and optimise.",
    ])

    add_content_slide(prs, "Option C: Recipe Optimisation", [
        "Scenario: optimise a bread baking recipe",
        "7 candidate factors: flour protein %, water ratio, yeast amount,",
        "  salt amount, kneading time, proof time, oven temperature",
        "3 responses: loaf volume (max), crumb score (max), crust color (target: 4)",
        "Budget: 25 total experimental runs",
        "",
        "Use doe-helper to design, simulate, analyse, and optimise.",
    ])

    add_content_slide(prs, "Capstone Requirements", [
        "Phase 1: Screening (doe init, doe generate, doe analyze)",
        "  Choose an appropriate screening design for all factors",
        "  Identify the 3-4 most important factors",
        "  Document your reasoning for dropping factors",
        "",
        "Phase 2: Characterisation & Optimisation",
        "  Run a CCD or Box-Behnken on surviving factors",
        "  Analyse with doe analyze, interpret the response surface",
        "  Optimise with doe optimize",
    ])

    add_content_slide(prs, "Capstone Deliverables", [
        "1. Completed config.json files for each phase",
        "2. doe report HTML output for screening and RSM phases",
        "3. Summary document (1-2 pages) covering:",
        "  Problem statement and factor selection rationale",
        "  Screening results and factor elimination decisions",
        "  RSM analysis and optimal settings",
        "  Practical recommendations",
        "4. Brief presentation (5 minutes) to the class",
    ])

    add_content_slide(prs, "Capstone Workflow with doe-helper", [
        "Step 1: doe init project/ --design plackett_burman",
        "Step 2: doe generate project/ && bash project/run.sh",
        "Step 3: doe analyze project/ -- identify significant factors",
        "Step 4: Create new config with reduced factors, type=ccd",
        "Step 5: doe generate project_rsm/ && bash project_rsm/run.sh",
        "Step 6: doe analyze project_rsm/ -- response surface analysis",
        "Step 7: doe optimize project_rsm/ -- find optimal settings",
        "Step 8: doe report project_rsm/ -- generate final report",
    ])

    add_exercise_slide(prs, "Capstone Project", [
        "1. Choose Option A, B, C, or propose your own scenario",
        "2. Complete the full workflow using doe-helper commands",
        "3. Produce all deliverables listed on the previous slide",
        "4. Prepare a 5-minute presentation covering:",
        "  What you learned about the system",
        "  Key factors and their effects",
        "  Optimal settings and predicted performance",
        "  What you would do next with more runs",
        "",
        "Time allocation: 90 minutes for the project, 30 minutes for presentations",
    ])

    add_section_divider(prs, "Course Summary", "Congratulations on completing the DOE Helper Training Course!")

    add_content_slide(prs, "What You've Learned", [
        "Module 1: Why DOE matters and the OVAT trap",
        "Module 2: doe-helper installation and configuration",
        "Module 3: Full factorial designs and interactions",
        "Module 4: Screening designs for many factors",
        "Module 5: Response surface methodology and optimisation",
        "Module 6: Analysis, interpretation, and reporting",
        "Module 7: Multi-response optimisation and advanced topics",
        "Module 8: Complete DOE workflow from problem to solution",
    ])

    add_content_slide(prs, "doe-helper Command Reference", [
        "doe init         Create experiment from template",
        "doe generate     Create design matrix and runner script",
        "doe record       Interactively record experimental results",
        "doe status       Show experiment progress",
        "doe info         Display design summary and evaluation metrics",
        "doe analyze      Run ANOVA, effects, diagnostics, plots",
        "doe optimize     Find optimal factor settings",
        "doe power        Compute statistical power",
        "doe augment      Extend an existing design",
        "doe report       Generate interactive HTML report",
        "doe export-worksheet  Export design as CSV or markdown",
    ])

    add_key_point_slide(prs, "Next Steps", [
        "Practice with the 221 built-in use cases",
        "Visit doehelper.com for documentation and examples",
        "Apply DOE to a real problem in your work",
        "Start small: 2-3 factors, full factorial",
        "Remember: a planned experiment beats trial-and-error every time",
    ])

    prs.save("slides/Module_08_Capstone.pptx")
    print("  Module 8 saved")


if __name__ == "__main__":
    print("Generating DOE Helper Training Course slides...")
    build_module_1()
    build_module_2()
    build_module_3()
    build_module_4()
    build_module_5()
    build_module_6()
    build_module_7()
    build_module_8()
    print("\nAll slides saved to training/slides/")
